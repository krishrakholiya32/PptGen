import os
import random
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree


def hex_to_rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip("#")
    if len(hex_color) == 3:
        hex_color = "".join(c * 2 for c in hex_color)
    r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


def is_dark(hex_color: str) -> bool:
    hex_color = hex_color.lstrip("#")
    r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    return (0.299 * r + 0.587 * g + 0.114 * b) < 128


def fit_image_to_slot(img_path: str, slot_w, slot_h, slide, left, top):
    from PIL import Image as PILImage
    try:
        with PILImage.open(img_path) as img:
            img_w, img_h = img.size

        slot_ratio = slot_w / slot_h
        img_ratio  = img_w / img_h

        # Add picture filling the full slot
        pic = slide.shapes.add_picture(img_path, left, top, slot_w, slot_h)

        # Only crop if ratio difference is significant
        if abs(img_ratio - slot_ratio) < 0.08:
            return pic

        if img_ratio > slot_ratio:
            # Image is wider than slot — crop sides
            visible_frac = slot_ratio / img_ratio
            excess = max(0.0, min(0.4, 1.0 - visible_frac))  # clamp to avoid over-cropping
            pic.crop_left  = excess / 2
            pic.crop_right = excess / 2
        else:
            # Image is taller than slot — crop top/bottom
            visible_frac = img_ratio / slot_ratio
            excess = max(0.0, min(0.4, 1.0 - visible_frac))  # clamp to avoid over-cropping
            pic.crop_top    = excess / 2
            pic.crop_bottom = excess / 2

        return pic
    except Exception as e:
        print(f"[Renderer] fit_image_to_slot failed: {e}")
        try:
            return slide.shapes.add_picture(img_path, left, top, slot_w, slot_h)
        except Exception:
            return None


THEMES = [
    {
        "name": "midnight",
        "bg": "1A1A2E", "surface": "16213E",
        "accent": "E94560", "accent2": "0F3460",
        "text": "EAEAEA", "text_muted": "A0A0B0",
        "font_title": "Arial Black", "font_body": "Calibri",
    },
    {
        "name": "forest",
        "bg": "1B2A1F", "surface": "243B2A",
        "accent": "4CAF7D", "accent2": "2E7D52",
        "text": "F0F7F1", "text_muted": "A8C4AE",
        "font_title": "Georgia", "font_body": "Calibri",
    },
    {
        "name": "coral",
        "bg": "FFFFFF", "surface": "FFF5F5",
        "accent": "F96167", "accent2": "2F3C7E",
        "text": "1A1A2E", "text_muted": "555577",
        "font_title": "Arial Black", "font_body": "Calibri",
    },
    {
        "name": "ocean",
        "bg": "021B2C", "surface": "042D4A",
        "accent": "00C9FF", "accent2": "0077B6",
        "text": "E8F4FD", "text_muted": "90C8E0",
        "font_title": "Trebuchet MS", "font_body": "Calibri",
    },
]


def _tint(hex_color: str, factor: float) -> str:
    """Lighten (factor>1) or darken (factor<1) a hex color."""
    h = hex_color.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    if factor > 1:
        r = int(r + (255 - r) * (factor - 1))
        g = int(g + (255 - g) * (factor - 1))
        b = int(b + (255 - b) * (factor - 1))
    else:
        r, g, b = int(r * factor), int(g * factor), int(b * factor)
    return f"{min(r,255):02X}{min(g,255):02X}{min(b,255):02X}"


def pick_theme(style: dict) -> dict:
    accent_raw = style.get("accent_color", "").lstrip("#").upper()
    bg_raw     = style.get("background_color", "").lstrip("#").upper()

    if not accent_raw or not bg_raw:
        return THEMES[0]

    dark = is_dark(f"#{bg_raw}")
    # Surface: slightly lighter/darker than bg for contrast in title band
    surface = _tint(bg_raw, 1.12) if dark else _tint(bg_raw, 0.92)
    text        = "EAEAEA" if dark else "1A1A2E"
    text_muted  = "A0A0B0" if dark else "666688"
    # accent2: slightly darker shade of accent for secondary elements
    accent2 = _tint(accent_raw, 0.75)

    return {
        "bg":         bg_raw,
        "surface":    surface,
        "accent":     accent_raw,
        "accent2":    accent2,
        "text":       text,
        "text_muted": text_muted,
        "font_title": style.get("primary_font", "Calibri"),
        "font_body":  "Calibri",
    }


IMAGE_VARIANTS_HEAVY = ["image_right", "image_left", "image_bottom"]
IMAGE_VARIANTS_LIGHT = ["image_right", "image_left", "image_top", "image_bottom", "full_image"]
TEXT_VARIANTS = ["title_content", "two_column", "accent_block", "centered"]


def pick_layout_variant(has_image: bool, bullet_count: int, recent: list) -> str:
    if has_image:
        pool = IMAGE_VARIANTS_HEAVY if bullet_count >= 4 else IMAGE_VARIANTS_LIGHT
    else:
        pool = TEXT_VARIANTS
    available = [v for v in pool if v not in recent[-2:]]
    if not available:
        available = pool
    return random.choice(available)


# ── Shape helpers ─────────────────────────────────────────────────────────────

def set_bg(slide, hex_color: str):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(hex_color)


def add_rect(slide, x, y, w, h, hex_fill):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(hex_fill)
    shape.line.fill.background()
    return shape


def add_textbox(slide, text, x, y, w, h,
                font_name="Calibri", font_size=18, bold=False,
                color="FFFFFF", align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    # Vertically center text within the box
    tf.vertical_anchor = 3   # middle  (1=top, 2=middle, 3=bottom — MSOPPTX uses 3 for middle)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = hex_to_rgb(color)
    return txBox


def add_textbox_multiline(slide, lines, x, y, w, h,
                          font_name="Calibri", font_size=16, bold=False,
                          color="FFFFFF", align=PP_ALIGN.LEFT,
                          line_spacing_pt=None):
    """Render multiple text lines into one textbox with even line spacing."""
    from pptx.util import Pt as _Pt
    from pptx.oxml.ns import qn as _qn
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing_pt:
            p.line_spacing = _Pt(line_spacing_pt)
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = _Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = hex_to_rgb(color)
    return txBox


def add_bullets(slide, bullets, x, y, w, h,
                font_name="Calibri", font_size=16,
                text_color="EAEAEA", accent_color="E94560"):
    """
    Place bullets filling the full slot height evenly.
    Distributes line_spacing and space_after dynamically based on bullet count
    so content expands to fill the available space — no empty gap at the bottom.
    """
    if not bullets:
        return
    bullets = [str(b)[:140] for b in bullets]
    n = len(bullets)

    # Convert slot height to points
    h_pt = h / 12700.0   # EMU → pt  (1 pt = 12700 EMU)

    # Each bullet has 2 text lines (one main + one wrap) at font_size
    # line_height = font_size * line_spacing_multiplier
    # space_after separates bullets
    # total = n * (2 * font_size * ls + sa) = h_pt
    # We solve for ls and sa with sa = ls * font_size * 0.8
    # n * (2 * font_size * ls + ls * font_size * 0.8) = h_pt
    # ls = h_pt / (n * font_size * 2.8)

    ls_multiplier = h_pt / (n * font_size * 2.8)
    # Clamp: min 1.0× (readable), max 3.0× (not ridiculous)
    ls_multiplier = max(1.0, min(3.0, ls_multiplier))
    line_spacing_pt = font_size * ls_multiplier
    space_after_pt  = font_size * ls_multiplier * 0.7

    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = Pt(line_spacing_pt)
        p.space_after   = Pt(space_after_pt)
        dot = p.add_run()
        dot.text = "● "
        dot.font.name = font_name
        dot.font.size = Pt(font_size - 1)
        dot.font.color.rgb = hex_to_rgb(accent_color)
        dot.font.bold = True
        run = p.add_run()
        run.text = str(bullet)
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.color.rgb = hex_to_rgb(text_color)


def _add_image_placeholder(slide, x, y, w, h, T):
    """Styled placeholder when no image is available."""
    add_rect(slide, x, y, w, h, T["accent2"])
    # Subtle diagonal stripes
    stripe_c = _tint(T["accent2"], 0.88)
    n_stripes = 6
    step = w // n_stripes
    for i in range(n_stripes):
        add_rect(slide, x + i * step, y, Pt(1.5), h, stripe_c)
    # Centered label
    add_textbox(slide, "[  IMAGE  ]",
                x, y + (h - Inches(0.3)) // 2, w, Inches(0.3),
                font_name=T["font_body"], font_size=11, bold=True,
                color=T["text_muted"], align=PP_ALIGN.CENTER)


def resolve_image(filename, image_dir):
    if not filename or not image_dir:
        return None
    path = os.path.join(image_dir, filename)
    return path if os.path.exists(path) else None


def get_bullets(data, max_bullets=5):
    bullets = data.get("bullets", [])
    if not isinstance(bullets, list):
        return []
    return bullets[:max_bullets]


def add_footer(slide, presentation_title: str, slide_num: int, W, H, T):
    add_rect(slide, 0, H - Inches(0.28), W, Inches(0.28), T["surface"])
    add_rect(slide, 0, H - Inches(0.28), W, Pt(1), T["accent"])
    add_textbox(slide, presentation_title[:60],
                Inches(0.4), H - Inches(0.26), W - Inches(1.5), Inches(0.24),
                font_name=T["font_body"], font_size=9, bold=False,
                color=T["text_muted"], align=PP_ALIGN.LEFT)
    add_textbox(slide, str(slide_num),
                W - Inches(0.8), H - Inches(0.26), Inches(0.6), Inches(0.24),
                font_name=T["font_body"], font_size=9, bold=True,
                color=T["accent"], align=PP_ALIGN.RIGHT)


def add_slide_header(slide, data, W, H, T, title_font_size=None):
    """Top accent bar + title band with word-wrap and auto font-size."""
    from pptx.enum.text import MSO_ANCHOR
    TITLE_BAND_H = Inches(1.1)
    title = data.get("title", "")
    if title_font_size is None:
        n = len(title)
        title_font_size = 28 if n <= 45 else (23 if n <= 70 else 19)
    add_rect(slide, 0, 0, W, Pt(6), T["accent"])
    add_rect(slide, 0, Pt(6), W, TITLE_BAND_H - Pt(6), T["surface"])
    txBox = slide.shapes.add_textbox(
        Inches(0.5), Pt(6), W - Inches(1.4), TITLE_BAND_H - Pt(6)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.name = T["font_title"]
    run.font.size = Pt(title_font_size)
    run.font.bold = True
    run.font.color.rgb = hex_to_rgb(T["text"])
    add_rect(slide, Inches(0.5), TITLE_BAND_H + Pt(2), Inches(4), Pt(2), T["accent"])


def add_slide_number_badge(slide, data, W, T):
    slide_num = data.get("slide_number", "")
    if slide_num:
        add_rect(slide, W - Inches(0.75), Inches(0.08), Inches(0.55), Inches(0.28), T["accent"])
        add_textbox(slide, str(slide_num),
                    W - Inches(0.75), Inches(0.08), Inches(0.55), Inches(0.28),
                    font_name=T["font_body"], font_size=10, bold=True,
                    color="FFFFFF", align=PP_ALIGN.CENTER)


# ── Slide Layouts ─────────────────────────────────────────────────────────────

# Shared geometry constants
TITLE_BAND_H  = Inches(1.1)   # height of the title area at top
FOOTER_H      = Inches(0.28)  # height of footer at bottom
PAD           = Inches(0.5)   # horizontal padding


def _content_top():    return TITLE_BAND_H + Inches(0.18)
def _content_bottom(H): return H - FOOTER_H - Inches(0.1)
def _content_h(H):     return _content_bottom(H) - _content_top()


def layout_title_slide(slide, data, W, H, T, presentation_title, slide_num):
    set_bg(slide, T["bg"])
    # Top half: accent fill gives the slide a bold, professional look
    split = H * 0.48
    add_rect(slide, 0, 0, W, split, T["accent"])
    # Thin accent2 divider at the split
    add_rect(slide, 0, split, W, Pt(5), T["accent2"])
    # Vertical accent bar on left
    add_rect(slide, 0, 0, Inches(0.22), H, T["accent2"])
    # Title — white text inside the accent half
    title = data.get("title", "Presentation")
    n = len(title)
    fsz = 46 if n <= 30 else (38 if n <= 50 else 30)
    add_textbox(slide, title,
                Inches(0.5), split * 0.18, W - Inches(1.0), split * 0.72,
                font_name=T["font_title"], font_size=fsz, bold=True,
                color="FFFFFF", align=PP_ALIGN.LEFT)
    # Subtitle — below the split in bg-coloured area
    subtitle = data.get("subtitle", "")
    sub_color = T["accent"] if not is_dark(f"#{T['bg']}") else T["accent"]
    if subtitle:
        lines = [l.strip() for l in subtitle.split("|") if l.strip()]
        if len(lines) > 1:
            add_textbox_multiline(slide, lines,
                Inches(0.5), split + Inches(0.35), W - Inches(1.0), Inches(1.1),
                font_name=T["font_body"], font_size=17, bold=False,
                color=sub_color, align=PP_ALIGN.LEFT, line_spacing_pt=24)
        else:
            add_textbox(slide, subtitle,
                        Inches(0.5), split + Inches(0.35), W - Inches(1.0), Inches(0.8),
                        font_name=T["font_body"], font_size=19, bold=False,
                        color=sub_color, align=PP_ALIGN.LEFT)
    # Short accent underline below subtitle
    add_rect(slide, Inches(0.5), split + Inches(1.4), Inches(3.5), Pt(3), T["accent"])
    add_footer(slide, presentation_title, slide_num, W, H, T)


def layout_image_right(slide, data, W, H, T, img_path, presentation_title, slide_num):
    set_bg(slide, T["bg"])
    img_x = W / 2 + Inches(0.1)
    img_w = W / 2 - Inches(0.5)
    img_y = _content_top()
    img_h = _content_h(H)
    if img_path:
        try:
            fit_image_to_slot(img_path, img_w, img_h, slide, img_x, img_y)
        except Exception:
            _add_image_placeholder(slide, img_x, img_y, img_w, img_h, T)
    else:
        _add_image_placeholder(slide, img_x, img_y, img_w, img_h, T)
    bullets = get_bullets(data)
    text_w = W / 2 - Inches(0.7)
    add_bullets(slide, bullets,
                PAD, _content_top(), text_w, _content_h(H),
                font_name=T["font_body"], font_size=16,
                text_color=T["text"], accent_color=T["accent"])
    add_slide_header(slide, data, W, H, T, title_font_size=26)
    add_slide_number_badge(slide, data, W, T)
    add_footer(slide, presentation_title, slide_num, W, H, T)


def layout_image_left(slide, data, W, H, T, img_path, presentation_title, slide_num):
    set_bg(slide, T["bg"])
    img_w = W / 2 - Inches(0.5)
    img_y = _content_top()
    img_h = _content_h(H)
    if img_path:
        try:
            fit_image_to_slot(img_path, img_w, img_h, slide, Inches(0.3), img_y)
        except Exception:
            _add_image_placeholder(slide, Inches(0.3), img_y, img_w, img_h, T)
    else:
        _add_image_placeholder(slide, Inches(0.3), img_y, img_w, img_h, T)
    text_x = W / 2 + Inches(0.2)
    text_w = W / 2 - Inches(0.7)
    bullets = get_bullets(data)
    add_slide_header(slide, data, W, H, T)
    add_bullets(slide, bullets,
                text_x, _content_top(), text_w, _content_h(H),
                font_name=T["font_body"], font_size=15,
                text_color=T["text"], accent_color=T["accent"])
    add_slide_number_badge(slide, data, W, T)
    add_footer(slide, presentation_title, slide_num, W, H, T)


def layout_image_top(slide, data, W, H, T, img_path, presentation_title, slide_num):
    set_bg(slide, T["bg"])
    img_h = H * 0.40
    if img_path:
        try:
            fit_image_to_slot(img_path, W, img_h, slide, 0, 0)
        except Exception:
            _add_image_placeholder(slide, 0, 0, W, img_h, T)
    else:
        _add_image_placeholder(slide, 0, 0, W, img_h, T)
    add_rect(slide, 0, img_h - Inches(0.65), W, Inches(0.65), "0D0D0D")
    add_rect(slide, 0, img_h, W, Pt(4), T["accent"])
    title = data.get("title", "")
    add_textbox(slide, title,
                Inches(0.5), img_h - Inches(0.62), W - Inches(1.0), Inches(0.58),
                font_name=T["font_title"], font_size=24, bold=True,
                color="FFFFFF", align=PP_ALIGN.LEFT)
    bullets = get_bullets(data, max_bullets=4)
    bul_top = img_h + Inches(0.2)
    bul_h   = H - bul_top - FOOTER_H - Inches(0.1)
    add_bullets(slide, bullets,
                PAD, bul_top, W - Inches(1.0), bul_h,
                font_name=T["font_body"], font_size=15,
                text_color=T["text"], accent_color=T["accent"])
    add_footer(slide, presentation_title, slide_num, W, H, T)


def layout_image_bottom(slide, data, W, H, T, img_path, presentation_title, slide_num):
    set_bg(slide, T["bg"])
    img_y = H * 0.54
    if img_path:
        try:
            fit_image_to_slot(img_path, W, H - img_y - FOOTER_H, slide, 0, img_y)
        except Exception:
            _add_image_placeholder(slide, 0, img_y, W, H - img_y - FOOTER_H, T)
    else:
        _add_image_placeholder(slide, 0, img_y, W, H - img_y - FOOTER_H, T)
    add_rect(slide, 0, img_y, W, Pt(4), T["accent"])
    add_slide_header(slide, data, W, H, T)
    add_slide_number_badge(slide, data, W, T)
    bul_top = _content_top()
    bul_h   = img_y - bul_top - Inches(0.1)
    bullets = get_bullets(data, max_bullets=3)
    add_bullets(slide, bullets,
                PAD, bul_top, W - Inches(1.0), bul_h,
                font_name=T["font_body"], font_size=16,
                text_color=T["text"], accent_color=T["accent"])
    add_footer(slide, presentation_title, slide_num, W, H, T)


def layout_full_image(slide, data, W, H, T, img_path, presentation_title, slide_num):
    set_bg(slide, T["bg"])
    if img_path:
        try:
            fit_image_to_slot(img_path, W, H, slide, 0, 0)
        except Exception:
            pass
    bullets = get_bullets(data, max_bullets=3)
    overlay_h = Inches(3.6) if bullets else Inches(2.0)
    add_rect(slide, 0, H - overlay_h, W, overlay_h, "0A0A0A")
    add_rect(slide, 0, H - overlay_h, W, Pt(3), T["accent"])
    title = data.get("title", "")
    title_y = H - overlay_h + Inches(0.18)
    add_textbox(slide, title,
                PAD, title_y, W - Inches(1.0), Inches(0.8),
                font_name=T["font_title"], font_size=30, bold=True,
                color="FFFFFF", align=PP_ALIGN.LEFT)
    if bullets:
        bul_top = title_y + Inches(0.85)
        bul_h   = H - bul_top - FOOTER_H - Inches(0.1)
        add_bullets(slide, bullets,
                    PAD, bul_top, W - Inches(1.0), bul_h,
                    font_name=T["font_body"], font_size=15,
                    text_color="EAEAEA", accent_color=T["accent"])
    add_footer(slide, presentation_title, slide_num, W, H, T)


def layout_title_content(slide, data, W, H, T, presentation_title, slide_num):
    set_bg(slide, T["bg"])
    bullets = get_bullets(data)
    add_bullets(slide, bullets,
                PAD, _content_top(), W - Inches(1.0), _content_h(H),
                font_name=T["font_body"], font_size=16,
                text_color=T["text"], accent_color=T["accent"])
    add_slide_header(slide, data, W, H, T)
    add_slide_number_badge(slide, data, W, T)
    add_footer(slide, presentation_title, slide_num, W, H, T)


def layout_two_column(slide, data, W, H, T, presentation_title, slide_num):
    set_bg(slide, T["bg"])
    bullets = get_bullets(data, max_bullets=8)
    mid = max(1, len(bullets) // 2)
    col_w = W / 2 - Inches(0.8)
    col_h = _content_h(H) - Inches(0.35)   # minus label strip
    label_y = _content_top()
    bul_y   = label_y + Inches(0.35)

    add_rect(slide, PAD, label_y, col_w, Inches(0.3), T["accent2"])
    add_textbox(slide, "Key Points",
                PAD, label_y, col_w, Inches(0.3),
                font_name=T["font_body"], font_size=11, bold=True,
                color="FFFFFF", align=PP_ALIGN.CENTER)
    add_bullets(slide, bullets[:mid],
                PAD, bul_y, col_w, col_h,
                font_name=T["font_body"], font_size=14,
                text_color=T["text"], accent_color=T["accent"])

    col2_x = W / 2 + Inches(0.3)
    add_rect(slide, col2_x, label_y, col_w, Inches(0.3), T["accent2"])
    add_textbox(slide, "Details",
                col2_x, label_y, col_w, Inches(0.3),
                font_name=T["font_body"], font_size=11, bold=True,
                color="FFFFFF", align=PP_ALIGN.CENTER)
    add_bullets(slide, bullets[mid:],
                col2_x, bul_y, col_w, col_h,
                font_name=T["font_body"], font_size=14,
                text_color=T["text"], accent_color=T["accent"])

    add_slide_header(slide, data, W, H, T)
    add_slide_number_badge(slide, data, W, T)
    add_footer(slide, presentation_title, slide_num, W, H, T)


def layout_accent_block(slide, data, W, H, T, presentation_title, slide_num):
    set_bg(slide, T["bg"])
    # Right accent bar with thin stripe pattern for depth
    bar_x = W - Inches(0.9)
    bar_w = Inches(0.9)
    add_rect(slide, bar_x, _content_top(), bar_w, _content_h(H), T["accent2"])
    add_rect(slide, bar_x + Inches(0.15), _content_top(), Pt(2), _content_h(H), T["accent"])
    add_rect(slide, bar_x + Inches(0.35), _content_top(), Pt(2), _content_h(H), T["accent"])
    bullets = get_bullets(data)
    add_bullets(slide, bullets,
                PAD, _content_top(), W - Inches(1.5), _content_h(H),
                font_name=T["font_body"], font_size=16,
                text_color=T["text"], accent_color=T["accent"])
    add_slide_header(slide, data, W, H, T)
    add_slide_number_badge(slide, data, W, T)
    add_footer(slide, presentation_title, slide_num, W, H, T)


def layout_centered(slide, data, W, H, T, presentation_title, slide_num):
    set_bg(slide, T["bg"])
    add_rect(slide, 0, 0, W, Pt(6), T["accent"])
    add_rect(slide, 0, H - FOOTER_H - Pt(6), W, Pt(6), T["accent"])
    title = data.get("title", "")
    # Title centered in its band
    add_textbox(slide, title,
                Inches(1.0), Inches(0.3), W - Inches(2.0), Inches(1.1),
                font_name=T["font_title"], font_size=34, bold=True,
                color=T["accent"], align=PP_ALIGN.CENTER)
    add_rect(slide, W / 2 - Inches(2), Inches(1.45), Inches(4), Pt(2), T["accent"])
    bullets = get_bullets(data)
    bul_top = Inches(1.6)
    bul_h   = H - bul_top - FOOTER_H - Inches(0.1)
    add_bullets(slide, bullets,
                Inches(1.5), bul_top, W - Inches(3.0), bul_h,
                font_name=T["font_body"], font_size=17,
                text_color=T["text"], accent_color=T["accent"])
    add_footer(slide, presentation_title, slide_num, W, H, T)


def layout_section_break(slide, data, W, H, T, presentation_title, slide_num):
    # Always use accent bg so section breaks stand out on both light and dark themes
    set_bg(slide, T["accent"])
    add_rect(slide, 0, 0, Inches(0.55), H, T["accent2"])
    add_rect(slide, Inches(0.55), 0, W - Inches(0.55), H * 0.45, T["accent2"])
    title = data.get("title", "")
    n = len(title)
    fsz = 44 if n <= 40 else (36 if n <= 60 else 28)
    add_textbox(slide, title,
                Inches(0.9), Inches(2.2), W - Inches(1.4), Inches(2.4),
                font_name=T["font_title"], font_size=fsz, bold=True,
                color="FFFFFF", align=PP_ALIGN.LEFT)
    subtitle = data.get("subtitle", "")
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.9), Inches(4.8), W - Inches(1.4), Inches(0.8),
                    font_name=T["font_body"], font_size=17, bold=False,
                    color="D0D8E8", align=PP_ALIGN.LEFT)
    add_footer(slide, presentation_title, slide_num, W, H, T)


# ── Main Renderer ─────────────────────────────────────────────────────────────

class PPTXRenderer:

    def render(self, plan: dict, style: dict, image_dir: str, output_path: str) -> str:
        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)
        W = prs.slide_width
        H = prs.slide_height

        T = pick_theme(style)
        slides_data = plan.get("slides", [])
        presentation_title = plan.get("title", "")
        recent_layouts = []

        for i, slide_data in enumerate(slides_data):
            ai_layout   = slide_data.get("layout_type", "title_content")
            slide       = prs.slides.add_slide(prs.slide_layouts[6])
            img_path    = resolve_image(slide_data.get("image"), image_dir)
            slide_num   = i + 1
            bullets     = get_bullets(slide_data)

            if ai_layout == "title_slide" or i == 0:
                layout_title_slide(slide, slide_data, W, H, T, presentation_title, slide_num)
                recent_layouts.append("title_slide")

            elif ai_layout == "section_break":
                layout_section_break(slide, slide_data, W, H, T, presentation_title, slide_num)
                recent_layouts.append("section_break")

            else:
                chosen = pick_layout_variant(
                    has_image=bool(img_path),
                    bullet_count=len(bullets),
                    recent=recent_layouts
                )
                recent_layouts.append(chosen)
                print(f"[Renderer] Slide {slide_num} '{slide_data.get('title','')[:30]}' -> {chosen}")

                dispatch = {
                    "image_right":   lambda s,d: layout_image_right(s, d, W, H, T, img_path, presentation_title, slide_num),
                    "image_left":    lambda s,d: layout_image_left(s, d, W, H, T, img_path, presentation_title, slide_num),
                    "image_top":     lambda s,d: layout_image_top(s, d, W, H, T, img_path, presentation_title, slide_num),
                    "image_bottom":  lambda s,d: layout_image_bottom(s, d, W, H, T, img_path, presentation_title, slide_num),
                    "full_image":    lambda s,d: layout_full_image(s, d, W, H, T, img_path, presentation_title, slide_num),
                    "two_column":    lambda s,d: layout_two_column(s, d, W, H, T, presentation_title, slide_num),
                    "accent_block":  lambda s,d: layout_accent_block(s, d, W, H, T, presentation_title, slide_num),
                    "centered":      lambda s,d: layout_centered(s, d, W, H, T, presentation_title, slide_num),
                    "title_content": lambda s,d: layout_title_content(s, d, W, H, T, presentation_title, slide_num),
                }
                dispatch.get(chosen, dispatch["title_content"])(slide, slide_data)
                pass  # transition hints removed (polluted speaker notes)

            recent_layouts = recent_layouts[-2:]

            if slide_data.get("speaker_notes"):
                try:
                    notes_tf = slide.notes_slide.notes_text_frame
                    existing = notes_tf.text.strip()
                    if slide_data["speaker_notes"] not in existing:
                        notes_tf.text = (f"{slide_data['speaker_notes']}\n\n{existing}"
                                         if existing else slide_data["speaker_notes"])
                except Exception:
                    pass

        prs.save(output_path)
        return output_path


def add_transition_hint(slide, layout_name: str):
    hints = {
        "image_right": "Transition: Fade", "image_left": "Transition: Push Right",
        "image_top": "Transition: Cover Down", "image_bottom": "Transition: Uncover Up",
        "full_image": "Transition: Dissolve", "two_column": "Transition: Wipe",
        "accent_block": "Transition: Fade", "centered": "Transition: Zoom",
        "title_content": "Transition: Fade",
    }
    hint = hints.get(layout_name, "Transition: Fade")
    try:
        notes_tf = slide.notes_slide.notes_text_frame
        existing = notes_tf.text.strip()
        notes_tf.text = f"{existing}\n\n[Design] {hint}" if existing else f"[Design] {hint}"
    except Exception:
        pass


renderer_pptx = PPTXRenderer()
