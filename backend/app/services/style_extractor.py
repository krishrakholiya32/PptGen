import os
import json
import logging
import subprocess
import collections
from pathlib import Path
from typing import Optional
from app.services.llm_service import llm

logger = logging.getLogger(__name__)


class StyleExtractor:
    """
    Extracts fonts, colors, layouts, and design intent from uploaded PPTX files.
    Returns a unified StyleProfile JSON used by the renderer.
    """

    async def extract(self, file_path: str) -> dict:
        ext = Path(file_path).suffix.lower()
        if ext == ".pptx":
            return await self._from_pptx(file_path)
        else:
            return self._default_style()

    async def _from_pptx(self, path: str) -> dict:
        from pptx import Presentation
        from pptx.util import Pt
        from pptx.dml.color import RGBColor

        prs = Presentation(path)
        fonts, bg_colors, text_colors, accent_colors = [], [], [], []

        # Extract from slide master first (most reliable source of theme)
        for layout in prs.slide_layouts:
            for ph in layout.placeholders:
                if ph.has_text_frame:
                    for para in ph.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.name:
                                fonts.append(run.font.name)

        # Extract from actual slides
        for slide in prs.slides:
            # Background color
            bg = slide.background.fill
            try:
                if bg.fore_color.rgb:
                    bg_colors.append(str(bg.fore_color.rgb))
            except Exception:
                pass

            for shape in slide.shapes:
                # Shape fill colors
                try:
                    fill = shape.fill
                    if fill.fore_color.rgb:
                        accent_colors.append(str(fill.fore_color.rgb))
                except Exception:
                    pass

                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.name:
                                fonts.append(run.font.name)
                            try:
                                if run.font.color.rgb:
                                    text_colors.append(str(run.font.color.rgb))
                            except Exception:
                                pass

        # Screenshot first slide for vision analysis
        layout_description = await self._vision_layout_analysis(path, "pptx")

        style = {
            "source_type": "pptx",
            "slide_width_emu": prs.slide_width,
            "slide_height_emu": prs.slide_height,
            "primary_font": collections.Counter(fonts).most_common(1)[0][0] if fonts else "Calibri",
            "fonts": list(set(fonts[:5])),
            "background_color": f"#{bg_colors[0]}" if bg_colors else "#FFFFFF",
            "text_color": f"#{collections.Counter(text_colors).most_common(1)[0][0]}" if text_colors else "#000000",
            "accent_color": f"#{collections.Counter(accent_colors).most_common(1)[0][0]}" if accent_colors else "#0070C0",
            "slide_count": len(prs.slides),
            "layout_description": layout_description
        }
        return style

    async def _vision_layout_analysis(self, file_path: str, file_type: str) -> str:
        """Use Groq vision model to analyze the visual layout of the first page/slide."""
        try:
            # Convert first page/slide to image
            img_path = await self._render_first_page(file_path, file_type)
            if not img_path:
                return "Standard layout"

            prompt = """Analyze this document page/slide layout. Describe:
1. Overall layout style (e.g. title-top, split-layout, full-bleed image, minimal text)
2. Color scheme (dark/light, primary colors used)
3. Typography feel (serif/sans-serif, large/small, bold headers)
4. Design style (corporate, modern, creative, minimal, colorful)
5. Content placement (where titles, images, text blocks are positioned)

Be concise — 3-4 sentences max."""

            return await llm.analyze_image(img_path, prompt)
        except Exception as e:
            logger.warning(f"Vision analysis failed for {file_path}: {e}")
            return "Standard layout"

    async def _render_first_page(self, file_path: str, file_type: str) -> Optional[str]:
        """Render first slide to PNG for vision analysis."""
        try:
            img_path = file_path.replace(f".{file_type}", "_preview.png")
            if file_type == "pptx":
                # Use an argument list (no shell=True) so the file path can never be
                # interpreted as shell syntax, regardless of its contents.
                result = subprocess.run(
                    ["libreoffice", "--headless", "--convert-to", "png", "--outdir", "/tmp", file_path],
                    capture_output=True,
                    timeout=60,
                )
                if result.returncode != 0:
                    logger.warning(
                        "LibreOffice conversion failed for %s (exit %s): %s",
                        file_path, result.returncode, result.stderr.decode(errors="replace").strip(),
                    )
                tmp = f"/tmp/{Path(file_path).stem}.png"
                if os.path.exists(tmp):
                    import shutil
                    shutil.copy(tmp, img_path)
            return img_path if os.path.exists(img_path) else None
        except subprocess.TimeoutExpired:
            logger.error("LibreOffice conversion timed out for %s", file_path)
            return None
        except FileNotFoundError:
            logger.error("LibreOffice binary not found — template vision analysis is disabled")
            return None
        except Exception:
            logger.exception("Unexpected error rendering preview for %s", file_path)
            return None

    def _default_style(self) -> dict:
        return {
            "source_type": "default",
            "primary_font": "Calibri",
            "fonts": ["Calibri"],
            "background_color": "#FFFFFF",
            "text_color": "#1A1A1A",
            "accent_color": "#0070C0",
            "layout_description": "Clean modern layout with title and content areas"
        }


extractor = StyleExtractor()
